#!/usr/bin/env python3

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image
import io

class PresentationCreator:
    def __init__(self, template_path):
        self.prs = Presentation(template_path)
        self.output_path = None

    def analyze_template(self):
        print(f"Презентация содержит {len(self.prs.slides)} слайдов")
        for i, slide in enumerate(self.prs.slides):
            print(f"\nСлайд {i+1}:")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if text:
                        print(f"  Текст: {text[:50]}...")
                if hasattr(shape, 'image'):
                    print(f"  Изображение найдено")

    def add_text_to_slide(self, slide_index, text, placeholder_index=0):
        if slide_index >= len(self.prs.slides):
            print(f"Слайд {slide_index+1} не существует")
            return

        slide = self.prs.slides[slide_index]
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame and i == placeholder_index:
                shape.text_frame.text = text
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                break

    def add_image_with_effects(self, slide_index, image_path, left, top, width=None, height=None):
        if slide_index >= len(self.prs.slides):
            print(f"Слайд {slide_index+1} не существует")
            return

        slide = self.prs.slides[slide_index]

        # Добавляем изображение
        if width and height:
            pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                          width=Inches(width), height=Inches(height))
        else:
            pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top))

        # Добавляем тень для объёма
        shadow = pic.shadow
        shadow.visible = True
        shadow.distance = Pt(4)
        shadow.direction = 45
        shadow.blur_radius = Pt(6)

        return pic

    def add_chart(self, slide_index, chart_type='column'):
        # Заготовка для добавления графиков
        pass

    def save_presentation(self, output_path):
        self.output_path = output_path
        self.prs.save(output_path)
        print(f"Презентация сохранена: {output_path}")

    def get_slide_layouts(self):
        print("Доступные макеты слайдов:")
        for i, layout in enumerate(self.prs.slide_layouts):
            print(f"  {i}: {layout.name}")

if __name__ == "__main__":
    # Создаем объект презентации
    creator = PresentationCreator("template.pptx")

    # Анализируем шаблон
    print("=== Анализ шаблона ===")
    creator.analyze_template()

    print("\n=== Доступные макеты ===")
    creator.get_slide_layouts()

    # Сохраняем как новый файл
    creator.save_presentation("company_presentation.pptx")